"""
PowerPoint document extractor for .pptx and .ppt files
"""
import time
import io
from typing import List
from .base_extractor import BaseExtractor, ExtractionResult


class PowerPointExtractor(BaseExtractor):
    """PowerPoint presentation text extraction"""
    
    @property
    def supported_mime_types(self) -> List[str]:
        return [
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',  # .pptx
            'application/vnd.ms-powerpoint',  # .ppt
            'application/mspowerpoint'
        ]
    
    @property
    def supported_extensions(self) -> List[str]:
        return ['pptx', 'ppt']
    
    def extract_text(self, file_data: bytes, filename: str = "") -> ExtractionResult:
        """
        Extract text from PowerPoint presentations
        
        Args:
            file_data: PowerPoint file bytes
            filename: Original filename
            
        Returns:
            ExtractionResult with extracted text and metadata
        """
        start_time = time.time()
        
        # Validate file size
        self.validate_file_size(file_data, max_size_mb=100)
        
        # Determine PowerPoint format
        is_pptx = filename.lower().endswith('.pptx') or b'PK' in file_data[:10]
        
        try:
            if is_pptx:
                return self._extract_pptx(file_data, filename, start_time)
            else:
                return self._extract_ppt(file_data, filename, start_time)
                
        except Exception as e:
            raise Exception(f"PowerPoint processing failed: {str(e)}")
    
    def _extract_pptx(self, file_data: bytes, filename: str, start_time: float) -> ExtractionResult:
        """Extract text from .pptx files using python-pptx"""
        try:
            from pptx import Presentation
            from pptx.exc import PackageNotFoundError
            
            # Load presentation from bytes
            prs = Presentation(io.BytesIO(file_data))
            
            all_text = ""
            slide_count = 0
            text_shapes_count = 0
            table_count = 0
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = f"\n=== Slide {slide_idx} ===\n"
                slide_has_content = False
                
                # Extract text from shapes
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'text') and shape.text.strip():
                            slide_text += shape.text.strip() + "\n"
                            text_shapes_count += 1
                            slide_has_content = True
                        
                        # Handle tables in slides
                        elif hasattr(shape, 'table'):
                            table = shape.table
                            table_text = f"\n[Slide {slide_idx} Table]\n"
                            table_has_data = False
                            
                            for row in table.rows:
                                row_data = []
                                row_has_data = False
                                
                                for cell in row.cells:
                                    cell_text = cell.text.strip()
                                    if cell_text:
                                        row_data.append(cell_text)
                                        row_has_data = True
                                    else:
                                        row_data.append("")
                                
                                if row_has_data:
                                    table_text += " | ".join(row_data) + "\n"
                                    table_has_data = True
                            
                            if table_has_data:
                                slide_text += table_text
                                table_count += 1
                                slide_has_content = True
                        
                    except Exception as shape_error:
                        print(f"    Warning: Could not extract from shape in slide {slide_idx}: {shape_error}")
                        continue
                
                # Extract slide notes if available
                try:
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        slide_text += f"\n[Slide {slide_idx} Notes]\n{notes_text}\n"
                        slide_has_content = True
                except Exception:
                    pass  # Notes extraction failed, continue
                
                if slide_has_content:
                    all_text += slide_text
                    slide_count += 1
            
            processing_time = time.time() - start_time
            
            # Extract presentation properties if available
            metadata = {
                'file_type': 'powerpoint_pptx',
                'filename': filename,
                'slides_processed': slide_count,
                'total_slides': len(prs.slides),
                'text_shapes': text_shapes_count,
                'tables': table_count,
                'processing_time_seconds': processing_time,
                'extractor': 'python-pptx'
            }
            
            try:
                core_props = prs.core_properties
                metadata.update({
                    'title': core_props.title or '',
                    'author': core_props.author or '',
                    'subject': core_props.subject or '',
                    'created': str(core_props.created) if core_props.created else '',
                    'modified': str(core_props.modified) if core_props.modified else ''
                })
            except Exception:
                pass  # Ignore metadata extraction errors
            
            print("PowerPoint (.pptx) processing completed:")
            print(f"  - Slides: {slide_count}/{len(prs.slides)}")
            print(f"  - Text shapes: {text_shapes_count}")
            print(f"  - Tables: {table_count}")
            print(f"  - Text length: {len(all_text)} characters")
            print(f"  - Processing time: {processing_time:.2f}s")
            
            return ExtractionResult(
                text=all_text,
                pages=slide_count,
                metadata=metadata,
                processing_time=processing_time
            )
            
        except ImportError:
            raise Exception("python-pptx not installed - cannot process .pptx files")
        except PackageNotFoundError:
            raise Exception("Invalid or corrupted PowerPoint presentation")
        except Exception as e:
            raise Exception(f"PPTX extraction failed: {str(e)}")
    
    def _extract_ppt(self, file_data: bytes, filename: str, start_time: float) -> ExtractionResult:
        """Handle .ppt files (legacy format)"""
        try:
            processing_time = time.time() - start_time
            
            # Legacy .ppt format is not well supported by Python libraries
            # Most libraries focus on the newer .pptx format
            
            metadata = {
                'file_type': 'powerpoint_ppt',
                'filename': filename,
                'processing_time_seconds': processing_time,
                'extractor': 'none',
                'note': 'Legacy .ppt format not supported - consider converting to .pptx'
            }
            
            print("PowerPoint (.ppt) processing:")
            print("  - Legacy .ppt format is not supported")
            print("  - Recommendation: Convert .ppt files to .pptx format for full support")
            print(f"  - Processing time: {processing_time:.2f}s")
            
            return ExtractionResult(
                text=f"[Unable to extract text from legacy .ppt file: {filename}]\n[Legacy .ppt format is not supported]\n[Recommendation: Please convert to .pptx format for full text extraction]",
                pages=1,
                metadata=metadata,
                processing_time=processing_time
            )
            
        except Exception as e:
            processing_time = time.time() - start_time
            
            metadata = {
                'file_type': 'powerpoint_ppt', 
                'filename': filename,
                'processing_time_seconds': processing_time,
                'extractor': 'none',
                'error': str(e),
                'note': 'Legacy .ppt format not supported'
            }
            
            return ExtractionResult(
                text=f"[Error processing .ppt file: {filename}]\n[Error: {str(e)}]\n[Recommendation: Convert to .pptx format]",
                pages=1,
                metadata=metadata,
                processing_time=processing_time
            )